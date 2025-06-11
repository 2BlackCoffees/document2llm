"""
@author Jean-Philippe Ulpiano
"""
import re
import json
from pprint import pformat, pprint
from typing import List, Dict, Tuple
from logging import Logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from domain.ichecker import DeckChecker, ArtisticSlideChecker, TextSlideChecker
from domain.allm_access import AbstractLLMAccess
from domain.llm_utils import LLMUtils
from domain.icontent_out import IContentOut
from domain.adocument2datastructure import ADocumentToDatastructure
from infrastructure.ppt_reader import PPTReader

class PowerPointToDataStructure(ADocumentToDatastructure): 
    TITLE_PARAMS: str = 'title_document'
    CHECKER_INSTANCE: str = 'checker_params'
    LLM_REQUEST_PARAMS: str = 'llm_request_params'
    DONE_TEXT: str = 'done_text'

    def __init__(self, document_path: str, slides_to_skip: List, slides_to_keep: List,\
                 logger: Logger, content_out: IContentOut, llm_utils: LLMUtils, \
                 selected_text_slide_requests: List, selected_artistic_slide_requests: List, \
                 selected_deck_requests: List, llm_access: AbstractLLMAccess, consider_bullets_for_crlf: bool= True):
        super().__init__(logger, content_out, llm_access)
        self.document =  Presentation(document_path)
        self.llm_utils = llm_utils
        self.document_path = document_path
        self.crlf_replacement = " * " if consider_bullets_for_crlf else " "
        self.slides_to_skip = slides_to_skip
        self.slides_to_keep = slides_to_keep
        self.selected_text_slide_requests = selected_text_slide_requests
        self.selected_artistic_slide_requests = selected_artistic_slide_requests
        self.selected_deck_requests = selected_deck_requests
        self.want_selected_text_slide_requests: bool = len(self.selected_text_slide_requests) > 0
        self.want_selected_artistic_slide_requests: bool = len(self.selected_artistic_slide_requests) > 0

    def __get_title_details(self, shape: Dict):
        self.logger.debug(f"Shape for title: {pformat(shape)}")
        title: str = ""
        if 'text' in shape['json']['shape']:
            title =  re.sub(r'^\s*', '', f"{shape['json']['shape']['text']}")      
            title =  re.sub(r'\n', '', title)      
            title =  re.sub(r'\s*$', '', title)      
        slide_info = f"Slide number {shape['json']['shape']['slide_number']} {title}"
        self.logger.info(slide_info)
        return title, slide_info

    def __get_slide_details(self, sorted_shapes: List, slide_number: int, shape_title: Dict):
            slide_shapes_content: List = [] 
            title: str = ''
            slide_info: str = ''
            reduced_slide_text: List = []
            pattern = re.compile("\S")
            title_found: bool = False
            if shape_title is not None:
                title, slide_info = self.__get_title_details(shape_title)
                title_found = True
                self.logger.debug(f'Found shape title: {pformat(shape_title)}')
            for shape in sorted_shapes:
                if (not title_found):
                    self.logger.debug(f'Searching if shape fits for a title: {pformat(shape)}')
                    if len(shape['raw_text']) > 0 and shape['json']['shape']['type'] in [str(MSO_SHAPE_TYPE.TEXT_BOX), str(MSO_SHAPE_TYPE.GROUP), str(MSO_SHAPE_TYPE.PLACEHOLDER)] and pattern.match(shape['raw_text']): 
                        title, slide_info = self.__get_title_details(shape)
                        shape['json']['shape']["is_title"] = "True"
                        title_found = True
                slide_shapes_content.append(shape['json']['shape'])
                reduced_slide_text.append(shape['raw_text'])
            if not title_found:
                self.logger.warning(f"No title found for slide: {slide_number}, see slide details below")
                self.logger.debug(pformat(sorted_shapes))
            return slide_shapes_content, title, slide_info, reduced_slide_text
    

    def __print_slide_keep_skip_info(self, keep_skip_info: str) -> None:
        self.logger.info(keep_skip_info)
        self.content_out.document(f"**{keep_skip_info}**")

    def _document_to_data_structure(self) -> List:

        data_structure: List = []
        deck_content: List = []
        slide_size: Tuple = (self.document.slide_width, self.document.slide_height)
        for slide_idx, slide in enumerate(self.document.slides):

            slide_number: int = slide_idx + 1
            if slide_number in self.slides_to_skip:
                self.__print_slide_keep_skip_info(f"Skipped slide number {slide_number} as per request.")
                continue

            if self.slides_to_keep is not None and len(self.slides_to_keep) > 0:
                if slide_number not in self.slides_to_keep:
                    continue
                else:
                    self.__print_slide_keep_skip_info(f"Keep slide number {slide_number} as per request.")
                

            if slide.element.get('show', '1' == '0'):
                self.__print_slide_keep_skip_info(f"Skipped hidden slide number {slide_number}.")
                continue
 
            self.logger.info(f"Analyzing slide number {slide_number}")
 
            shape_descriptions_text_only: list = [] 
            shape_descriptions_with_graphics: list = [] 
            for shape in slide.shapes:
                if shape.has_text_frame:
                    PPTReader.add_text_box_info(slide_number, shape, False, shape_descriptions_text_only, slide_size)
                    if self.want_selected_artistic_slide_requests:
                        PPTReader.add_text_box_info(slide_number, shape, True, shape_descriptions_with_graphics, slide_size)
 
                elif shape.has_table: 
                    table = shape.table
                    table_md_str: str = ""
                    first_row: bool = True
                    for row in table.rows:
                        table_md_str += "\n|"
                        for cell in row.cells:
                            if cell.text_frame is not None:
                              text: str = cell.text_frame.text.replace("\n", " ")
                              table_md_str += text + "|"
                        if first_row: table_md_str += "\n" + ("-" * max(3, len(table_md_str)))
                        first_row = False

                    PPTReader.add_table_info(slide_number, shape, table, table_md_str, False, shape_descriptions_text_only, slide_size)
                    if self.want_selected_artistic_slide_requests:
                        PPTReader.add_table_info(slide_number, shape, table, table_md_str, True, shape_descriptions_with_graphics, slide_size)

                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP: 
                    PPTReader.add_group_info(slide_number, shape, False, shape_descriptions_text_only, slide_size)
                    if self.want_selected_artistic_slide_requests:
                        PPTReader.add_group_info(slide_number, shape, True, shape_descriptions_with_graphics, slide_size)
               
                else: 
                    PPTReader.add_shape_type_info(slide_number, shape, False, shape_descriptions_text_only, slide_size)
                    if self.want_selected_artistic_slide_requests:
                        PPTReader.add_shape_type_info(slide_number, shape, True, shape_descriptions_with_graphics, slide_size)
 
            sorted_shapes: List = PPTReader.get_sorted_shapes_by_pos_y(shape_descriptions_text_only)
            title_value: str = None
            shape_title: Dict = None
            title_found: bool = False
            if hasattr(slide.shapes, "title") and hasattr(slide.shapes.title, "text") and \
               slide.shapes.title.text is not None and len(slide.shapes.title.text) > 0:
                title_value = slide.shapes.title.text
                title_value: str = re.sub(r'\*\*', '', title_value)
                self.logger.debug(f"Title shape found contains text: {title_value}")
                for shape_description in sorted_shapes:
                    raw_text: str = re.sub(r'\n$', '', shape_description["raw_text"])
                    raw_text: str = re.sub(r'\*\*', '', raw_text)
                    if raw_text == title_value:
                        shape_title = shape_description
                        shape_title["json"]["shape"]["is_title"] = True
                        title_found = True
                        self.logger.debug(f'Found a shape with same text as title:\n{json.dumps(shape_description, sort_keys=True, indent=2, separators=(",", ": "))}')
                if not title_found:
                    PPTReader.add_created_title(slide_number, title_value, False, shape_descriptions_text_only, slide_size)
                    if self.want_selected_artistic_slide_requests:
                        PPTReader.add_created_title(slide_number, shape, True, shape_descriptions_with_graphics, slide_size)
                    shape_title = shape_description
                    shape_descriptions_text_only.append(shape_description)

            self.logger.debug(f'Not sorted shapes: {json.dumps(shape_descriptions_text_only, sort_keys=True, indent=2, separators=(",", ": "))}')
            self.logger.debug(f'Sorted shapes: {json.dumps(sorted_shapes, sort_keys=True, indent=2, separators=(",", ": "))}')
            self.logger.debug(f'shape_title: {json.dumps(shape_title, sort_keys=True, indent=2, separators=(",", ": "))}')

            slide_shapes_content, title, slide_info, reduced_slide_text = self.__get_slide_details(sorted_shapes, slide_number, shape_title)
            slide_content: Dict = {
                "slide_info": slide_info,
                "title": title,
                "shapes": slide_shapes_content,
                "reduced_slide_text": reduced_slide_text
            }

            slide_info: str = f'{slide_number} {title}'
            if self.want_selected_text_slide_requests or self.want_selected_artistic_slide_requests:
                data_structure.append(
                    {
                        self.TITLE_PARAMS: (1, f"Analyzing slide {slide_info}"),
                        self.CHECKER_INSTANCE: None,
                        self.LLM_REQUEST_PARAMS: None,
                        self.DONE_TEXT: None
                    }
                )

                if self.want_selected_text_slide_requests:
                    data_structure.append(
                        {
                            self.TITLE_PARAMS: (2, f"Check of text content for slide {slide_number}"),
                            self.CHECKER_INSTANCE: TextSlideChecker(self.llm_utils, self.selected_text_slide_requests, f' (Slide {slide_idx + 1})', f' (Slide {slide_idx + 1})'),
                            self.LLM_REQUEST_PARAMS: (slide_content["shapes"].copy(), False, slide_info),
                            self.DONE_TEXT: f"Text slide request {slide_number} {title}"
                        }
                    )

                if self.want_selected_artistic_slide_requests:
                    data_structure.append(
                        {
                            self.TITLE_PARAMS: (2, f"Check of artistic content for slide {slide_number}"),
                            self.CHECKER_INSTANCE: ArtisticSlideChecker(self.llm_utils, self.selected_artistic_slide_requests, f' (Slide {slide_idx + 1})', f' (Slide {slide_idx + 1})'),
                            self.LLM_REQUEST_PARAMS: (slide_content["shapes"].copy(), False, slide_info),
                            self.DONE_TEXT: f"Artistic slide request {slide_number} {title}"
                        }
                    )
                   
            deck_content.append(slide_content)

        if len(self.selected_deck_requests) > 0:
            formatted_deck_content_list: List = [ f'Slide {slide_number + 1}, {slide_content["title"]}:\n{json.dumps(slide_content["reduced_slide_text"])}' \
                                                  for slide_number, slide_content in enumerate(deck_content) ]
            data_structure.append(
                {
                    self.TITLE_PARAMS: (1, f"Check of text content and flow for the whole deck"),
                    self.CHECKER_INSTANCE: DeckChecker(self.llm_utils, self.selected_deck_requests, f' (Deck)', f' (Deck)'),
                    self.LLM_REQUEST_PARAMS: (formatted_deck_content_list, True, 'Whole deck'),
                    self.DONE_TEXT: f"Full deck request"
                }
            )
        return data_structure

    def _get_title_rank_title_str_as_tuple(self, data: Dict) -> Tuple:
        return data[self.TITLE_PARAMS]

    def _get_checker_instance(self, data: Dict) -> any:
        return data[self.CHECKER_INSTANCE]

    def _get_llm_parameters_requests_as_tuple(self, data: any) -> Tuple:
        return data[self.LLM_REQUEST_PARAMS]
    
    def _get_done_text(self, data: any) -> str:
        return data[self.DONE_TEXT]



        