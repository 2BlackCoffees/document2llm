import re
import json
import traceback
from pprint import pformat
from typing import List, Dict
from logging import Logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from domain.ichecker import IChecker, DeckChecker, ArtisticSlideChecker, TextSlideChecker
from domain.allm_access import AbstractLLMAccess
from domain.llm_utils import LLMUtils
from domain.ppt_reader import PPTReader
from domain.icontent_out import IContentOut

class PPT2GPT: 
    def __init__(self, document_path: str, slides_to_skip: List, \
                 logger: Logger, content_out: IContentOut, llm_utils: LLMUtils, \
                 selected_text_slide_requests: List, selected_artistic_slide_requests: List, \
                 selected_deck_requests: List, llm_access: AbstractLLMAccess, consider_bullets_for_crlf: bool= True):
        self.content_out = content_out
        self.logger = logger
        self.document =  Presentation(document_path)
        self.llm_utils = llm_utils
        self.document_path = document_path
        self.crlf_replacement = " * " if consider_bullets_for_crlf else " "
        self.slides_to_skip = slides_to_skip
        self.selected_text_slide_requests = selected_text_slide_requests
        self.selected_artistic_slide_requests = selected_artistic_slide_requests
        self.selected_deck_requests = selected_deck_requests
        self.llm_access = llm_access

        try:
            self.__ppt2gpt()
        except Exception as err:                    
            self.logger.warning(f"Caught exception {err=}\n {type(err)=}\n {traceback.print_exc()}\n Leaving application.")
        self.content_out.flush_and_close()

    def __get_title_details(self, shape: Dict):
        self.logger.debug(f"Shape for title: {pformat(shape)}")
        title =  re.sub(r'^\s*', '', f"{shape['json']['shape']['text']}")      
        title =  re.sub(r'\n', '', title)      
        title =  re.sub(r'\s*$', '', title)      
        slide_info = f"Slide number {shape['json']['shape']['slide_number']} {title}"
        self.logger.info(slide_info)
        return title, slide_info

    def __get_slide_details(self, sorted_shapes: List, slide_number: int, shape_title: Dict):
            slide_shapes_content: List = [] # self.utils.get_llm_instructions() + "\n"
            title: str = ''
            slide_info: str = ''
            reduced_slide_text: List = []
            pattern = re.compile("\S")
            title_found: bool = False
            if shape_title is not None:
                title, slide_info = self.__get_title_details(shape_title)
                title_found = True
                self.logger.debug(f'Found shape title: {pformat(shape_title)}')
            for shape in sorted_shapes:
                if (not title_found):
                    self.logger.debug(f'Searching if shape fits for a title: {pformat(shape)}')
                    if len(shape['raw_text']) > 0 and shape['json']['shape']['type'] in [str(MSO_SHAPE_TYPE.TEXT_BOX), str(MSO_SHAPE_TYPE.GROUP), str(MSO_SHAPE_TYPE.PLACEHOLDER)] and pattern.match(shape['raw_text']): 
                        title, slide_info = self.__get_title_details(shape)
                        shape['json']['shape']["is_title"] = "True"
                        title_found = True
                slide_shapes_content.append(shape['json']['shape'])
                reduced_slide_text.append(shape['raw_text'])
            if not title_found:
                self.logger.warning(f"No title found for slide: {slide_number}, see slide details below")
                self.logger.debug(pformat(sorted_shapes))
            return slide_shapes_content, title, slide_info, reduced_slide_text
    
    def __send_llm_requests_and_expand_output(self, content_to_check: List, print_title: bool) -> None:

            result = self.llm_access.check(content_to_check)

            for response in result:
                if print_title: 
                    self.content_out.add_title(2, response['request_name'])
                else:
                    self.content_out.document(f"**{response['request_name']}**")

                self.content_out.document(response['response'])
            
    def __ppt2gpt(self):
        deck_content: List = []
        for slide_idx, slide in enumerate(self.document.slides):

            slide_number: int = slide_idx + 1
            if slide_number in self.slides_to_skip:
                self.logger.info(f"Skipping slide number {slide_number} as per request.")
                self.content_out.add_title(1, f"Skipped slide number {slide_number} as per request")
                continue
            # TODO: Does not work, will have to be fixed
            if slide.element.get('show', '1' == '0'):
                self.logger.info(f'Skipping hidden slide number {slide_number}')
                self.content_out.add_title(1, f"Skipped hidden slide number {slide_number}")
                continue
 
            self.logger.info(f"Analyzing slide number {slide_number}")
 
            shape_descriptions: list = [] 
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_descriptions.append(PPTReader.get_text_box_info(slide_number, shape))
 
                elif shape.has_table: 
                    table = shape.table
                    table_elements: list = []
                    table_str: str = ""
                    for row_idx, row in enumerate(table.rows, start = 1):
                        table_str += "\n|"
                        for col_idx, cell in enumerate(row.cells, start=1):
                            if cell.text_frame is not None:
                              text: str = cell.text_frame.text.replace("\n", " ")
                              table_elements.append({'row': row_idx, 'col': col_idx, 'text': text})
                              table_str += text + "|"
 
                    shape_descriptions.append(PPTReader.get_table_info(slide_number, shape, table, table_elements, table_str))
 
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP: 
                    shape_descriptions.append(PPTReader.get_group_info(slide_number, shape))
               
                else: 
                    shape_descriptions.append(PPTReader.get_shape_type_info(slide_number, shape))
 
            title_value: str = None
            shape_title: Dict = None
            title_found: bool = False
            if hasattr(slide.shapes, "title") and hasattr(slide.shapes.title, "text") and slide.shapes.title.text is not None and len(slide.shapes.title.text) > 0:
                title_value = slide.shapes.title.text
                for shape_description in shape_descriptions:
                    if shape_description["raw_text"] == title_value:
                        shape_title = shape_description
                        shape_title["json"]["shape"]["is_title"] = True
                        title_found = True
                if not title_found:
                    shape_description: Dict = PPTReader.create_title(slide_number, title_value)
                    shape_title = shape_description
                    shape_descriptions.append(shape_description)

            sorted_shapes: List = PPTReader.get_sorted_shapes_by_pos_y(shape_descriptions)

            slide_shapes_content, title, slide_info, reduced_slide_text = self.__get_slide_details(sorted_shapes, slide_number, shape_title)
            slide_content: Dict = {
                "slide_info": slide_info,
                "title": title,
                "shapes": slide_shapes_content,
                "reduced_slide_text": reduced_slide_text
            }
            #title = re.sub(r'[^\w\d]', ' ', title)
            self.content_out.add_title(1, f"Analyzing slide {slide_number} {title}")

            if len(self.selected_text_slide_requests) > 0:
                self.content_out.add_title(2, f"Check of text content for slide {slide_number}")
                checker: IChecker = TextSlideChecker(self.llm_utils, self.selected_text_slide_requests, f' (Slide {slide_idx + 1})', f' (Slide {slide_idx + 1})')
                self.llm_access.set_checker(checker)
                self.__send_llm_requests_and_expand_output(slide_content["shapes"], False)

            if len(self.selected_artistic_slide_requests) > 0:
                self.content_out.add_title(2, f"Check of artistic content for slide {slide_number}")
                checker: IChecker = ArtisticSlideChecker(self.llm_utils, self.selected_artistic_slide_requests, f' (Slide {slide_idx + 1})', f' (Slide {slide_idx + 1})')
                self.llm_access.set_checker(checker)
                self.__send_llm_requests_and_expand_output(slide_content["shapes"], False)
            deck_content.append(slide_content)

        if len(self.selected_deck_requests) > 0:
            self.content_out.add_title(1, f"Check of text content and flow for the whole deck")
            formatted_deck_content_list: List = [ f'Slide {slide_number + 1}, {slide_content["title"]}:\n{json.dumps(slide_content["reduced_slide_text"])}' \
                                                  for slide_number, slide_content in enumerate(deck_content) ]
            checker: IChecker = DeckChecker(self.llm_utils, self.selected_deck_requests, f' (Deck)', f' (Deck)')
            self.llm_access.set_checker(checker)
            self.__send_llm_requests_and_expand_output(formatted_deck_content_list, True)

        