from typing import Dict, List
import re
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint, pformat
from pathlib import Path
class PPTReader:

    @staticmethod
    def is_paragraph(text: str, paragraph_start_min_word_length: str = 3, paragraph_start_min_word_numbers: str = 1) -> bool:
        minwords: int = int(paragraph_start_min_word_numbers)  
        regexp: str = f'(\\w{{{paragraph_start_min_word_length},}}\\b\\s+){{{minwords},}}\w{{{paragraph_start_min_word_length},}}\\b'
        paragraph_found: bool = (re.search(regexp, text) is not None)
        return paragraph_found

    @staticmethod  
    def __check_recursively_for_text(shape, text_list):
        for cur_shape in shape.shapes:
            if cur_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                PPTReader.__check_recursively_for_text(cur_shape, text_list)
            else:
                if hasattr(cur_shape, "text"):
                    text_list.append(cur_shape.text)
        return text_list
    
    def _get_shape_graphical_infos(shape: Dict, shape_details: Dict) -> Dict:


        if hasattr(shape, 'left') and hasattr(shape, 'right'): 
            shape_details["shape"]["position"] = {
                    "from_left": shape.left,
                    "from_top": shape.top
            }

        if hasattr(shape, 'width') and hasattr(shape, 'height'): 
            shape_details["shape"]["size"] = {
                    "width (EMU)": shape.width,
                    "height (EMU)": shape.height
            }

        if hasattr(shape, 'rotation'): 
            shape_details["shape"]["rotation_degrees"] = shape.rotation,
        shape_details["shape"]["shape_fore_color"] = "Transparent"

        if hasattr(shape, 'fill'):
            if shape.fill.type == MSO_FILL.SOLID:
                if shape.fill.fore_color.type == MSO_COLOR_TYPE.SCHEME:
                    shape_details["shape"]["shape_fore_color"] = f'Scheme color: {shape.fill.fore_color.theme_color}'
                elif shape.fill.fore_color.type == MSO_COLOR_TYPE.RGB:
                    shape_details["shape"]["shape_fore_color"] = f'(R, G, B): {shape.fill.fore_color.rgb}'
                else:
                    shape_details["shape"]["shape_fore_color"] = "Transparent"

        if hasattr(shape, 'line'):
            line_style: Dict = {}
            line_fill = shape.line.fill
            if line_fill.type == MSO_FILL.SOLID:
                line_color = line_fill.fore_color
                if line_color.type == MSO_COLOR_TYPE.SCHEME:
                    line_style["line_color"] =  f'Scheme color: {line_color.theme_color}'
                elif line_color.type == MSO_COLOR_TYPE.RGB:
                    line_style["line_color"] = f'Hexadecimal:0xRRGGBB: 0x{line_color.rgb}'
                else:
                    line_style["line_color"] = "Transparent"

            if hasattr(shape.line, "width"):
                line_style["line_width_points"] = shape.line.width.pt
            shape_details["shape"]["line"] = line_style

        if hasattr(shape, 'text_frame'):
            text_style: List = []
            for paragraph in shape.text_frame.paragraphs:
                font_details: Dict = {}
                for run in paragraph.runs:
                    font_details["font_name"] = run.font.name
                    font_details["font_size (PT)"] = run.font.size.pt if run.font.size is not None else "Not defined"
                    font_details["text_impacted"] = run.text
                    if run.font.color.type == MSO_COLOR_TYPE.SCHEME:
                        font_details["text_color"] =  f'Scheme color: {run.font.color.theme_color}'
                    elif run.font.color.type == MSO_COLOR_TYPE.RGB:
                        font_details["text_color"] = f'Hexadecimal:0xRRGGBB: 0x{run.font.color.rgb}'
                    else:
                        font_details["text_color"] = "None"           
                    text_style.append(font_details)    
            shape_details["shape"]["font_details"] = text_style
            shape_details["shape"]["text"] = shape.text_frame.text

        return shape_details

    def _get_shape_infos(slide_number: int, shape: Dict, need_graphical: bool) -> Dict:
        shape_details: Dict =  {
            "shape": {
                "slide_number": slide_number,

            }
        }
        if hasattr(shape, 'text_frame'):
            shape_details["shape"]["text"] = shape.text_frame.text
            
        if need_graphical:
            shape_details = PPTReader._get_shape_graphical_infos(shape, shape_details)
            
        return shape_details

    def __encapsulate_shape(shape: Dict, text: str, json_shape: Dict) -> Dict:
        new_json = json_shape.copy()
        if hasattr(shape, 'shape_type'):
            new_json ["shape"]["type"] = str(shape.shape_type)
        elif 'shape_type' in shape:
            new_json ["shape"]["type"] = shape['shape_type']
        else:
            new_json ["shape"]["type"] = 'NA'

        return {
            'y': shape.top if hasattr(shape, 'top') else shape['top'], 
            'x': shape.left if hasattr(shape, 'left') else shape['left'], 
            'json': new_json,
            'raw_text': str(new_json["shape"]["text"]) if text is None else text
        }
    def add_text_box_info(slide_number: int, shape: Dict, need_graphical: bool, shape_descriptions: List) -> None:
        json_shape: Dict = PPTReader._get_shape_infos(slide_number, shape, need_graphical)
        json_shape["shape"]["type"] = str(MSO_SHAPE_TYPE.TEXT_BOX)
        json_shape["shape"]["is_title"] = "False"
        final_json: Dict = PPTReader.__encapsulate_shape(shape, None, json_shape)
        if PPTReader.is_paragraph(final_json['raw_text']):
            shape_descriptions.append(final_json)
    
    def add_group_info(slide_number: int, shape: Dict, need_graphical: bool, shape_descriptions: List) -> None:
        json_shape: Dict = PPTReader._get_shape_infos(slide_number, shape, need_graphical)
        text_list: list = []
        text: str = "\n".join(PPTReader.__check_recursively_for_text(shape, text_list))
        final_json: Dict = PPTReader.__encapsulate_shape(shape, text, json_shape)
        if PPTReader.is_paragraph(final_json['raw_text']):
            shape_descriptions.append(final_json)
    
    def add_table_info(slide_number: int, shape: Dict, table: Dict, table_str: str, need_graphical: bool, shape_descriptions: List) -> None:
        json_shape: Dict = PPTReader._get_shape_infos(slide_number, shape, need_graphical)
        json_shape["shape"]["table_size"] = {
             "number_cols": len(table.rows[0].cells),
             "number_rowss": len(table.rows)
        }
        json_shape["shape"]["mytype"] = "table"
        json_shape["shape"]["table_cells"] = table_str # [ element for element in table_elements]

        shape_descriptions.append(PPTReader.__encapsulate_shape(shape, table_str, json_shape))

    def add_shape_type_info(slide_number: int, shape: Dict, need_graphical: bool, shape_descriptions: List) -> None:
        json_shape: Dict = PPTReader._get_shape_infos(slide_number, shape, need_graphical)
        json_shape["shape"]["type"] = "shape_type"
        text: str = ""
        if "text" in json_shape["shape"].keys():
            text = json_shape["shape"]["text"]

        shape_descriptions.append(PPTReader.__encapsulate_shape(shape, text, json_shape))

    def add_created_title(slide_number: int, title_value: str, need_graphical: bool, shape_descriptions: List) -> None:
        shape: Dict = {
            'top': 0,
            'left': 0,
            'right': 0,            
            'width': len(title_value),
            'height': 10,
            'shape_type': "ForcedTitle"
        }
        json_shape: Dict = PPTReader._get_shape_infos(slide_number, shape, need_graphical)
        json_shape["shape"]["type"] = str(MSO_SHAPE_TYPE.TEXT_BOX)
        json_shape["shape"]["is_title"] = "True"
        json_shape["shape"]["text"] = title_value
        shape_descriptions.append(PPTReader.__encapsulate_shape(shape, title_value, json_shape))


    def get_sorted_shapes_by_pos_y(shapes: List) -> List:
        return sorted(shapes, key = lambda shape_dict: shape_dict['y'])

